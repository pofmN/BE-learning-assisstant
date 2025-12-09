"""
Document text extraction service.
Supports PDF, DOCX, PPTX, and TXT files.
"""
import io
import logging
from typing import Optional
import magic
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


class DocumentExtractor:
    """Extract text content from various document formats."""
    
    SUPPORTED_MIME_TYPES = {
        "application/pdf": "PDF",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "DOCX",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "PPTX",
        "text/plain": "TXT",
    }
    
    def extract_text(self, file_bytes: bytes, filename: Optional[str] = None) -> str:
        """
        Extract text from document bytes.
        
        Args:
            file_bytes: Raw file content as bytes
            filename: Optional filename for logging
            
        Returns:
            Extracted text content
            
        Raises:
            ValueError: If file type is unsupported
        """
        mime_type = magic.from_buffer(file_bytes, mime=True)
        
        if mime_type not in self.SUPPORTED_MIME_TYPES:
            raise ValueError(
                f"Unsupported file type: {mime_type}. "
                f"Supported types: {', '.join(self.SUPPORTED_MIME_TYPES.values())}"
            )
        
        logger.info(f"Extracting text from {self.SUPPORTED_MIME_TYPES[mime_type]} file")
        
        if mime_type == "application/pdf":
            return self._extract_pdf(file_bytes)
        elif "wordprocessingml" in mime_type:
            return self._extract_docx(file_bytes)
        elif "presentationml" in mime_type:
            return self._extract_pptx(file_bytes)
        elif mime_type.startswith("text/"):
            return self._extract_text(file_bytes)
        
        raise ValueError(f"Unsupported MIME type: {mime_type}")
    
    def _extract_pdf(self, file_bytes: bytes) -> str:
        """Extract text from PDF."""
        try:
            pdf = PdfReader(io.BytesIO(file_bytes))
            text_parts = []
            
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(page_text)
                    
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            raise ValueError(f"Failed to extract PDF: {str(e)}")
    
    def _extract_docx(self, file_bytes: bytes) -> str:
        """Extract text from DOCX."""
        try:
            doc = DocxDocument(io.BytesIO(file_bytes))
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            raise ValueError(f"Failed to extract DOCX: {str(e)}")
    
    def _extract_pptx(self, file_bytes: bytes) -> str:
        """Extract text from PPTX."""
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
                    
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            raise ValueError(f"Failed to extract PPTX: {str(e)}")
    
    def _extract_text(self, file_bytes: bytes) -> str:
        """Extract text from plain text file."""
        try:
            return file_bytes.decode("utf-8", errors="ignore")
        except Exception as e:
            logger.error(f"Text extraction error: {e}")
            raise ValueError(f"Failed to decode text file: {str(e)}")