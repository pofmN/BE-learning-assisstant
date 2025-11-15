"""
Document parsing utilities for PDF, DOCX, and PPTX files.
"""
import os
from typing import Optional

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from PDF file.

    Args:
        file_path: Path to PDF file

    Returns:
        Extracted text content
    """
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        raise ValueError(f"Error extracting text from PDF: {str(e)}")


def extract_text_from_docx(file_path: str) -> str:
    """
    Extract text from DOCX file.

    Args:
        file_path: Path to DOCX file

    Returns:
        Extracted text content
    """
    try:
        doc = DocxDocument(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        raise ValueError(f"Error extracting text from DOCX: {str(e)}")


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from PPTX file.

    Args:
        file_path: Path to PPTX file

    Returns:
        Extracted text content
    """
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        raise ValueError(f"Error extracting text from PPTX: {str(e)}")


def extract_text_from_document(file_path: str, file_type: str) -> str:
    """
    Extract text from document based on file type.

    Args:
        file_path: Path to document file
        file_type: Type of file (pdf, docx, pptx)

    Returns:
        Extracted text content

    Raises:
        ValueError: If file type is not supported or extraction fails
    """
    if not os.path.exists(file_path):
        raise ValueError(f"File not found: {file_path}")

    file_type = file_type.lower()

    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type == "docx":
        return extract_text_from_docx(file_path)
    elif file_type == "pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
